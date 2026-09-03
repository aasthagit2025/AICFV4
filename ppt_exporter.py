from __future__ import annotations

from io import BytesIO
from typing import Iterable

import pandas as pd
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt


COLORS = {
    "ink": RGBColor(42, 45, 58),
    "muted": RGBColor(105, 111, 125),
    "line": RGBColor(224, 228, 235),
    "accent": RGBColor(230, 65, 82),
    "green": RGBColor(64, 148, 96),
    "amber": RGBColor(218, 151, 48),
    "red": RGBColor(191, 74, 74),
    "blue": RGBColor(64, 112, 180),
    "pale": RGBColor(247, 248, 251),
}


def clean(value: object) -> str:
    if pd.isna(value):
        return ""
    return " ".join(str(value).replace("\n", " ").split()).strip()


def shorten(value: object, limit: int = 145) -> str:
    text = clean(value)
    if len(text) <= limit:
        return text
    return text[: limit - 3].rstrip() + "..."


def confidence_counts(report: pd.DataFrame) -> dict[str, int]:
    if "ici_classification" in report.columns:
        levels = ["High Confidence", "Medium Confidence", "Low Confidence", "Not Scorable"]
        counts = report.get("ici_classification", pd.Series(dtype=str)).value_counts()
        return {level: int(counts.get(level, 0)) for level in levels}
    levels = ["High Confidence", "Medium Confidence", "Low Confidence", "Not Scorable"]
    counts = report.get("confidence_level", pd.Series(dtype=str)).value_counts()
    return {level: int(counts.get(level, 0)) for level in levels}


def add_text(slide, text: str, x: float, y: float, w: float, h: float, size: int = 18,
             bold: bool = False, color: RGBColor | None = None, align=None):
    box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    frame = box.text_frame
    frame.clear()
    frame.word_wrap = True
    paragraph = frame.paragraphs[0]
    paragraph.text = text
    if align is not None:
        paragraph.alignment = align
    for run in paragraph.runs:
        run.font.name = "Calibri"
        run.font.size = Pt(size)
        run.font.bold = bold
        run.font.color.rgb = color or COLORS["ink"]
    return box


def add_title(slide, title: str, subtitle: str = ""):
    add_text(slide, title, 0.55, 0.35, 12.1, 0.62, size=28, bold=True)
    if subtitle:
        add_text(slide, subtitle, 0.58, 1.0, 11.6, 0.38, size=13, color=COLORS["muted"])
    line = slide.shapes.add_shape(1, Inches(0.58), Inches(1.46), Inches(12.1), Inches(0.02))
    line.fill.solid()
    line.fill.fore_color.rgb = COLORS["line"]
    line.line.color.rgb = COLORS["line"]


def add_footer(slide, page_label: str):
    add_text(slide, "AI Insight Confidence Framework", 0.58, 7.05, 4.8, 0.22, size=9, color=COLORS["muted"])
    add_text(slide, page_label, 11.4, 7.05, 1.2, 0.22, size=9, color=COLORS["muted"], align=PP_ALIGN.RIGHT)


def add_metric(slide, label: str, value: str, x: float, y: float, color: RGBColor):
    shape = slide.shapes.add_shape(1, Inches(x), Inches(y), Inches(2.85), Inches(1.0))
    shape.fill.solid()
    shape.fill.fore_color.rgb = COLORS["pale"]
    shape.line.color.rgb = COLORS["line"]
    add_text(slide, value, x + 0.15, y + 0.13, 2.55, 0.35, size=22, bold=True, color=color)
    add_text(slide, label, x + 0.15, y + 0.55, 2.55, 0.25, size=10, color=COLORS["muted"])


def bullet_lines(slide, rows: Iterable[dict], x: float, y: float, w: float, line_h: float, limit: int = 5):
    current_y = y
    rows = list(rows)
    if not rows:
        add_text(slide, "No rows met this condition in the current report.", x, current_y, w, 0.35, size=14, color=COLORS["muted"])
        return
    for idx, row in enumerate(rows):
        if idx >= limit:
            break
        theme = shorten(row.get("theme", ""), 44)
        insight = shorten(row.get("insight_text", ""), 128)
        score = clean(row.get("ici_score", row.get("weighted_score", "")))
        prefix = f"{theme}"
        if score:
            prefix += f" | ICI {score}"
        add_text(slide, prefix, x, current_y, w, 0.22, size=11, bold=True, color=COLORS["accent"])
        add_text(slide, insight, x, current_y + 0.25, w, 0.43, size=12)
        current_y += line_h


def sorted_rows(report: pd.DataFrame, mask, ascending: bool, limit: int = 6) -> list[dict]:
    data = report.loc[mask].copy()
    sort_col = "ici_score" if "ici_score" in data.columns else "weighted_score"
    if sort_col in data.columns:
        data[sort_col] = pd.to_numeric(data[sort_col], errors="coerce")
        data = data.sort_values(sort_col, ascending=ascending)
    return data.head(limit).to_dict("records")


def add_confidence_slide(prs: Presentation, report: pd.DataFrame):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_title(slide, "ICI separates insights by confidence", "Every scored insight still requires human researcher sign-off.")
    counts = confidence_counts(report)
    total = max(1, sum(counts.values()))
    colors = {
        "High Confidence": COLORS["green"],
        "Medium Confidence": COLORS["amber"],
        "Low Confidence": COLORS["red"],
        "Not Scorable": COLORS["muted"],
    }
    y = 2.0
    for level, count in counts.items():
        width = 7.2 * count / total
        add_text(slide, level, 0.75, y - 0.02, 2.45, 0.28, size=12, bold=True)
        bar = slide.shapes.add_shape(1, Inches(3.15), Inches(y), Inches(max(0.15, width)), Inches(0.27))
        bar.fill.solid()
        bar.fill.fore_color.rgb = colors[level]
        bar.line.color.rgb = colors[level]
        add_text(slide, str(count), 10.55, y - 0.02, 1.0, 0.28, size=12, bold=True, color=colors[level])
        y += 0.72
    add_footer(slide, "Confidence view")


def add_dimension_explanation_slide(prs: Presentation, report: pd.DataFrame):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_title(slide, "Each AICF score explains what was strong or weak", "These plain-language notes help managers understand why an insight was marked ready, check recommended, or review required.")

    explanation_cols = [
        ("Evidence", "evidence_strength_explanation"),
        ("Triangulation", "triangulation_explanation"),
        ("Business relevance", "business_relevance_explanation"),
        ("Actionability", "actionability_explanation"),
        ("Bias control", "bias_risk_explanation"),
    ]
    usable = report.copy()
    if "weighted_score" in usable.columns:
        usable["weighted_score"] = pd.to_numeric(usable["weighted_score"], errors="coerce")
        usable = usable.sort_values("weighted_score", ascending=False)
    row = usable.iloc[0].to_dict() if not usable.empty else {}
    add_text(slide, shorten(row.get("theme", "Example insight"), 72), 0.72, 1.72, 11.6, 0.28, size=13, bold=True, color=COLORS["accent"])
    add_text(slide, shorten(row.get("insight_text", ""), 210), 0.72, 2.05, 11.6, 0.55, size=12)

    y = 2.95
    for label, col in explanation_cols:
        add_text(slide, label, 0.78, y, 2.15, 0.25, size=11, bold=True, color=COLORS["accent"])
        add_text(slide, shorten(row.get(col, "No explanation available."), 150), 2.55, y, 9.6, 0.38, size=11)
        y += 0.68
    add_footer(slide, "Why it scored this way")


def build_pptx_report(report: pd.DataFrame, title: str = "AICF Insight Confidence Report", subtitle: str = "") -> bytes:
    report = report.copy()
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    counts = confidence_counts(report)
    status_series = report.get("review_status", pd.Series(dtype=str))
    high_count = counts.get("High Confidence", 0)
    medium_count = counts.get("Medium Confidence", 0)
    low_count = counts.get("Low Confidence", 0)
    not_scorable_count = counts.get("Not Scorable", 0)

    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_text(slide, title, 0.7, 0.75, 11.8, 0.8, size=34, bold=True)
    add_text(slide, subtitle or "Confidence-rated synthesis; human researcher sign-off is required", 0.72, 1.55, 10.7, 0.4, size=16, color=COLORS["muted"])
    add_metric(slide, "Total scored insights", str(len(report)), 0.8, 2.45, COLORS["ink"])
    add_metric(slide, "High Confidence", str(high_count), 3.95, 2.45, COLORS["green"])
    add_metric(slide, "Medium / Low", str(medium_count + low_count), 7.1, 2.45, COLORS["amber"])
    add_metric(slide, "Not Scorable", str(not_scorable_count), 10.25, 2.45, COLORS["red"])
    add_text(slide, "AICF creates a confidence-rated evidence trail. High, Medium, and Low Confidence are decision aids only; every finding requires named human researcher sign-off.", 0.85, 4.05, 11.5, 0.75, size=18)
    add_footer(slide, "Overview")

    add_confidence_slide(prs, report)
    add_dimension_explanation_slide(prs, report)

    ready_rows = sorted_rows(
        report,
        report.get("review_status", pd.Series(dtype=str)).eq("High Confidence"),
        ascending=False,
        limit=6,
    )
    ready_title = "High Confidence findings have the strongest evidence"
    ready_subtitle = "These findings still require named human researcher sign-off before use."
    if not ready_rows:
        ready_title = "Most supported findings still need researcher review"
        ready_subtitle = "PPT-only reviews can identify strong claims, but source tables should be checked before client use."
        ready_rows = sorted_rows(
            report,
            report.get("ici_classification", report.get("confidence_level", pd.Series(dtype=str))).isin(["High Confidence", "Medium Confidence"]),
            ascending=False,
            limit=6,
        )
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_title(slide, ready_title, ready_subtitle)
    bullet_lines(slide, ready_rows, 0.75, 1.85, 11.6, 0.82, limit=6)
    add_footer(slide, "Ready insights")

    review_rows = sorted_rows(
        report,
        report.get("review_status", pd.Series(dtype=str)).isin(["Medium Confidence", "Low Confidence", "Not Scorable"]),
        ascending=True,
        limit=6,
    )
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_title(slide, "Review flags show where human judgment adds value", "AICF does not block these insights; it identifies where evidence, wording, or interpretation should be checked.")
    bullet_lines(slide, review_rows, 0.75, 1.85, 11.6, 0.82, limit=6)
    add_footer(slide, "Review watch-outs")

    evidence_rows = report[report.get("evidence_note", pd.Series(dtype=str)).astype(str).str.len() > 20].head(4).to_dict("records")
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_title(slide, "Evidence notes make each finding auditable", "The deck keeps a traceable link between the insight, the table or survey signal, and the confidence rating.")
    y = 1.85
    for row in evidence_rows:
        add_text(slide, shorten(row.get("theme", ""), 54), 0.75, y, 11.6, 0.22, size=11, bold=True, color=COLORS["accent"])
        add_text(slide, shorten(row.get("evidence_note", ""), 175), 0.75, y + 0.27, 11.6, 0.46, size=12)
        y += 0.95
    add_footer(slide, "Evidence examples")

    story = report.loc[report.get("theme", pd.Series(dtype=str)).eq("Complete Story")]
    summary = report.loc[report.get("theme", pd.Series(dtype=str)).eq("Overall Summary")]
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_title(slide, "The final story connects evidence to action", "Summary and story rows help researchers move from scored outputs to client narrative.")
    story_text = story.iloc[-1]["insight_text"] if not story.empty else ""
    summary_text = summary.iloc[-1]["insight_text"] if not summary.empty else ""
    if not summary_text:
        summary_text = (
            f"AICF scored {len(report)} extracted insights. "
            f"{high_count} are High Confidence, {medium_count} are Medium Confidence, {low_count} are Low Confidence, and {not_scorable_count} are Not Scorable. Human sign-off is required for all scored findings."
        )
    if not story_text:
        story_text = (
            "The overall story should be finalized by the researcher after reviewing the scored rows, "
            "especially where PowerPoint chart values need validation against source tables."
        )
    add_text(slide, shorten(summary_text, 430), 0.78, 1.85, 11.6, 1.45, size=15)
    add_text(slide, shorten(story_text, 520), 0.78, 3.65, 11.6, 1.85, size=15, color=COLORS["ink"])
    add_footer(slide, "Story")

    appendix_rows = report.head(6).to_dict("records")
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_title(slide, "Appendix sample shows the scoring trail", "Representative rows from the detailed report.")
    y = 1.75
    for row in appendix_rows:
        line = f"{shorten(row.get('theme', ''), 34)} | ICI {clean(row.get('ici_score', ''))} | {clean(row.get('ici_classification', row.get('review_status', '')))}"
        add_text(slide, line, 0.75, y, 11.7, 0.22, size=10, bold=True, color=COLORS["accent"])
        add_text(slide, shorten(row.get("insight_text", ""), 140), 0.75, y + 0.23, 11.7, 0.35, size=10)
        y += 0.74
    add_footer(slide, "Appendix")

    output = BytesIO()
    prs.save(output)
    return output.getvalue()
