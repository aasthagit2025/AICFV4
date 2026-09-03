"""Extract comparable metrics from uploaded tables, charts and analysis output,
then bind each metric to the individual insight it supports or contradicts.

The v5 banner-table parser only recognised tables containing the literal strings
``"frequency row"`` and ``"vertical % row"``. Those markers come from one
tabulation package's export layout. Every other banner table, every Excel
crosstab and the project's own sample file produced zero metric records, so the
cross-source layer silently found nothing and flagged every comparative claim as
unsupported.

v6 keeps that layout as a fast path and adds a general parser that recognises the
ordinary shape of a market research table: a header row of column labels, an
optional base row, and labelled data rows carrying two or more numeric cells.
Everything the parser did or failed to do is reported back so the researcher can
see why a claim was not matched.
"""

from __future__ import annotations

import io
import re
from dataclasses import dataclass, field
from typing import Any

import pandas as pd

from aicf_config import STOPWORDS

MAX_SHEETS = 20
MAX_ROWS_PER_SHEET = 4000


# ---------------------------------------------------------------------------
# Text and number helpers
# ---------------------------------------------------------------------------

def compact_text(value: object, limit: int = 360) -> str:
    text = " ".join(str(value or "").replace("\n", " ").split()).strip()
    return text if len(text) <= limit else text[: limit - 3].rstrip() + "..."


def normalize_compare_text(value: object) -> str:
    text = str(value or "").lower().replace("&", " and ")
    return " ".join(re.sub(r"[^a-z0-9]+", " ", text).split())


def tokenize(text: object) -> set[str]:
    return {
        token
        for token in re.findall(r"[a-zA-Z][a-zA-Z0-9-]{2,}", str(text).lower())
        if token not in STOPWORDS
    }


def question_codes(text: object) -> set[str]:
    return set(re.findall(r"\b[a-z]{1,4}\d+[a-z0-9_]*\b", str(text or "").lower()))


def value_as_number(value: object) -> float | None:
    if value is None:
        return None
    if isinstance(value, bool):
        return None
    if isinstance(value, (int, float)):
        return None if pd.isna(value) else float(value)
    text = str(value).strip().replace(",", "")
    if not text or text.lower() in {"nan", "-", "na", "n/a", "*"}:
        return None
    percent = "%" in text
    text = text.replace("%", "").strip()
    if not re.fullmatch(r"-?\d+(\.\d+)?", text):
        return None
    number = float(text)
    return number / 100 if percent else number


def format_metric_value(value: float) -> str:
    return f"{value * 100:.1f}%" if abs(value) <= 1.000001 else f"{value:,.1f}"


def file_bytes(uploaded_file) -> bytes:
    if uploaded_file is None:
        return b""
    if hasattr(uploaded_file, "getvalue"):
        return uploaded_file.getvalue()
    if hasattr(uploaded_file, "seek"):
        uploaded_file.seek(0)
    return uploaded_file.read()


# ---------------------------------------------------------------------------
# Diagnostics
# ---------------------------------------------------------------------------

@dataclass
class ParseDiagnostics:
    """What the parser saw in each file. Surfaced in the app's Diagnostics tab."""

    entries: list[dict[str, Any]] = field(default_factory=list)

    def add(self, source: str, sheet: str, outcome: str, detail: str = "") -> None:
        self.entries.append({"file": source, "sheet": sheet, "outcome": outcome, "detail": detail})

    def to_frame(self) -> pd.DataFrame:
        return pd.DataFrame(self.entries or [{"file": "-", "sheet": "-", "outcome": "No files parsed", "detail": ""}])


# ---------------------------------------------------------------------------
# Raw sheet reading
# ---------------------------------------------------------------------------

def read_raw_sheets(uploaded_file) -> list[tuple[str, pd.DataFrame]]:
    """Read every sheet of a workbook with no header assumption."""
    name = str(getattr(uploaded_file, "name", "uploaded file"))
    data = file_bytes(uploaded_file)
    if not data:
        return []
    lowered = name.lower()
    try:
        if lowered.endswith((".csv", ".txt", ".tsv")):
            sep = "\t" if lowered.endswith(".tsv") else ","
            for encoding in ("utf-8-sig", "latin1"):
                try:
                    frame = pd.read_csv(
                        io.BytesIO(data), header=None, encoding=encoding, sep=sep,
                        dtype=object, on_bad_lines="skip", nrows=MAX_ROWS_PER_SHEET,
                    )
                    return [(name, frame)]
                except UnicodeDecodeError:
                    continue
            return []
        if lowered.endswith((".xlsx", ".xlsm", ".xls")):
            excel = pd.ExcelFile(io.BytesIO(data))
            return [
                (sheet, pd.read_excel(excel, sheet_name=sheet, header=None, dtype=object, nrows=MAX_ROWS_PER_SHEET))
                for sheet in excel.sheet_names[:MAX_SHEETS]
            ]
    except Exception:
        return []
    return []


def nonempty_cells(row: pd.Series) -> list[tuple[int, str]]:
    cells = []
    for idx, value in enumerate(row.tolist()):
        text = str(value if value is not None else "").replace("\n", " ").strip()
        if text and text.lower() != "nan":
            cells.append((idx, text))
    return cells


# ---------------------------------------------------------------------------
# Generic banner-table parser
# ---------------------------------------------------------------------------

BASE_ROW_CUES = ["base", "total respondents", "sample size", "unweighted", "weighted base", "n ="]
TITLE_CUES = ["table", "q", "question"]


def _is_label(text: str) -> bool:
    return value_as_number(text) is None and len(text.strip()) >= 2 and not re.fullmatch(r"\([A-Z0-9]+\)", text.strip())


def _find_header_row(table: pd.DataFrame, data_row: int, lookback: int = 8) -> int | None:
    """The nearest preceding row that reads like a set of column labels."""
    best, best_score = None, 0
    for idx in range(max(0, data_row - lookback), data_row):
        cells = nonempty_cells(table.iloc[idx])
        labels = [(col, text) for col, text in cells if _is_label(text)]
        numbers = sum(1 for _, text in cells if value_as_number(text) is not None)
        if len(labels) < 2 or numbers > len(labels):
            continue
        score = len(labels) - numbers + (2 if idx == data_row - 1 else 0)
        if score > best_score:
            best, best_score = idx, score
    return best


def _column_label(table: pd.DataFrame, header_idx: int, col_idx: int) -> tuple[str, str]:
    """Return (label, parent group) for a column, merging a banner group row above."""
    label = str(table.iat[header_idx, col_idx] if col_idx < table.shape[1] else "").replace("\n", " ").strip()
    if label.lower() == "nan" or re.fullmatch(r"\([A-Z0-9]+\)", label):
        label = ""
    parent = ""
    if header_idx - 1 >= 0:
        for probe in range(col_idx, -1, -1):
            candidate = str(table.iat[header_idx - 1, probe] or "").replace("\n", " ").strip()
            if candidate and candidate.lower() != "nan" and not re.fullmatch(r"\([A-Z0-9]+\)", candidate):
                parent = candidate
                break
    return label, parent


def _nearest_title(table: pd.DataFrame, row_idx: int, lookback: int = 12) -> str:
    for idx in range(row_idx - 1, max(-1, row_idx - lookback) - 1, -1):
        cells = nonempty_cells(table.iloc[idx])
        if len(cells) != 1:
            continue
        text = cells[0][1]
        if len(text) > 6 and value_as_number(text) is None:
            return compact_text(text, 200)
    return ""


def parse_generic_table(source: str, sheet: str, raw: pd.DataFrame,
                        diagnostics: ParseDiagnostics | None = None) -> list[dict[str, Any]]:
    """Extract one record per labelled data row carrying two or more numbers."""
    table = raw.fillna("")
    if table.empty or table.shape[1] < 2:
        return []

    records: list[dict[str, Any]] = []
    current_base = ""

    for row_idx in range(len(table)):
        cells = nonempty_cells(table.iloc[row_idx])
        if len(cells) < 3:
            continue
        row_label = cells[0][1] if _is_label(cells[0][1]) else ""
        if not row_label:
            continue

        lowered = row_label.lower()
        numeric_cells = [(col, value_as_number(text)) for col, text in cells[1:]]
        numeric_cells = [(col, value) for col, value in numeric_cells if value is not None]
        if len(numeric_cells) < 2:
            continue

        if any(cue in lowered for cue in BASE_ROW_CUES):
            current_base = f"{row_label}: " + ", ".join(f"{format_metric_value(v)}" for _, v in numeric_cells[:6])
            continue

        header_idx = _find_header_row(table, row_idx)
        if header_idx is None:
            continue

        values = []
        for col, value in numeric_cells:
            label, parent = _column_label(table, header_idx, col)
            if not label or normalize_compare_text(label) in {"total", "total respondents", "all"}:
                continue
            values.append({"label": label, "parent": parent, "value": value})
        if len(values) < 2:
            continue

        records.append({
            "source_file": source,
            "sheet": sheet,
            "title": _nearest_title(table, header_idx) or sheet,
            "metric": row_label,
            "base_note": current_base,
            "values": values,
        })

    if diagnostics is not None:
        if records:
            diagnostics.add(source, sheet, f"{len(records)} comparable metric row(s) found",
                            "; ".join(r["metric"] for r in records[:5]))
        else:
            diagnostics.add(source, sheet, "No comparable metric rows found",
                            "Needs a header row of column labels and data rows with 2+ numbers.")
    return records


def parse_tabhouse_table(source: str, sheet: str, raw: pd.DataFrame) -> list[dict[str, Any]]:
    """Fast path for exports that mark frequency and percentage rows explicitly."""
    table = raw.fillna("")
    text_blob = " ".join(str(v).lower() for v in table.head(60).values.ravel())
    if "frequency row" not in text_blob:
        return []

    records: list[dict[str, Any]] = []
    for row_idx in range(len(table)):
        row_text = " ".join(text for _, text in nonempty_cells(table.iloc[row_idx])).lower()
        if "frequency row" not in row_text or "filter frequency row" in row_text:
            continue
        metric_label = next((t for _, t in nonempty_cells(table.iloc[row_idx]) if "frequency row" not in t.lower()), "")
        if not metric_label:
            continue
        pct_idx = None
        for probe in range(row_idx + 1, min(row_idx + 4, len(table))):
            if "vertical %" in " ".join(t for _, t in nonempty_cells(table.iloc[probe])).lower():
                pct_idx = probe
                break
        if pct_idx is None:
            continue
        header_idx = _find_header_row(table, row_idx)
        if header_idx is None:
            continue
        values = []
        for col in range(table.shape[1]):
            value = value_as_number(table.iloc[pct_idx].iloc[col])
            if value is None:
                continue
            label, parent = _column_label(table, header_idx, col)
            if not label or normalize_compare_text(label) in {"total", "total respondents"}:
                continue
            values.append({"label": label, "parent": parent, "value": value})
        if len(values) >= 2:
            records.append({
                "source_file": source, "sheet": sheet, "title": _nearest_title(table, header_idx) or sheet,
                "metric": metric_label, "base_note": "", "values": values,
            })
    return records


def extract_ppt_chart_metrics(uploaded_file, diagnostics: ParseDiagnostics | None = None) -> list[dict[str, Any]]:
    from pptx import Presentation

    source = str(getattr(uploaded_file, "name", "PowerPoint report"))
    data = file_bytes(uploaded_file)
    if not data:
        return []
    try:
        presentation = Presentation(io.BytesIO(data))
    except Exception as exc:
        if diagnostics is not None:
            diagnostics.add(source, "-", "Could not open PowerPoint file", str(exc))
        return []

    records: list[dict[str, Any]] = []
    image_only_slides = 0
    for slide_index, slide in enumerate(presentation.slides, start=1):
        title = next(
            (compact_text(shape.text, 120) for shape in slide.shapes
             if str(getattr(shape, "text", "")).strip() and 2 <= len(str(shape.text).split()) <= 18),
            f"Slide {slide_index}",
        )
        found_chart = False
        for shape in slide.shapes:
            if not getattr(shape, "has_chart", False):
                continue
            found_chart = True
            chart = shape.chart
            try:
                categories = [str(c).replace("\n", " ").strip() for c in chart.plots[0].categories]
            except Exception:
                continue
            for series in chart.series:
                try:
                    series_name = str(getattr(series, "name", "") or "Chart value").strip()
                    values = list(series.values)
                except Exception:
                    continue
                points = [
                    {"label": label, "parent": "PowerPoint chart", "value": value_as_number(value)}
                    for label, value in zip(categories, values)
                    if label and value_as_number(value) is not None
                ]
                if len(points) >= 2:
                    records.append({
                        "source_file": source, "sheet": f"Slide {slide_index}",
                        "title": compact_text(f"Slide {slide_index} {title}", 200),
                        "metric": series_name, "base_note": "", "values": points,
                    })
        if not found_chart and any(shape.shape_type == 13 for shape in slide.shapes):
            image_only_slides += 1

    if diagnostics is not None:
        diagnostics.add(source, "PowerPoint charts", f"{len(records)} readable chart series",
                        f"{image_only_slides} slide(s) contain pasted images whose values cannot be read; "
                        "upload the source table for those." if image_only_slides else "")
    return records


def extract_metric_records(uploaded_files, diagnostics: ParseDiagnostics | None = None) -> list[dict[str, Any]]:
    records: list[dict[str, Any]] = []
    for uploaded_file in uploaded_files or []:
        source = str(getattr(uploaded_file, "name", "uploaded file"))
        if source.lower().endswith(".pptx"):
            records.extend(extract_ppt_chart_metrics(uploaded_file, diagnostics))
            continue
        sheets = read_raw_sheets(uploaded_file)
        if not sheets and diagnostics is not None:
            diagnostics.add(source, "-", "Could not read this file as a table", "Check the file format and encoding.")
        for sheet, raw in sheets:
            fast = parse_tabhouse_table(source, sheet, raw)
            if fast:
                records.extend(fast)
                if diagnostics is not None:
                    diagnostics.add(source, sheet, f"{len(fast)} metric row(s) found (tabulation export layout)", "")
                continue
            records.extend(parse_generic_table(source, sheet, raw, diagnostics))
    return records


# ---------------------------------------------------------------------------
# Binding a metric record to one insight
# ---------------------------------------------------------------------------

DECISIVE_HIGH = ["most", "highest", "strongest", "best", "winner", "winning", "leads", "outperforms", "top"]
DECISIVE_LOW = ["lowest", "least", "weakest", "worst", "trails", "lags"]


def claim_direction(text: str) -> str | None:
    lowered = str(text or "").lower()
    if any(cue in lowered for cue in DECISIVE_LOW):
        return "low"
    if any(cue in lowered for cue in DECISIVE_HIGH):
        return "high"
    return None


def _relevance(row_text: str, record: dict[str, Any]) -> int:
    record_text = f"{record.get('title', '')} {record.get('metric', '')}"
    row_codes, record_codes = question_codes(row_text), question_codes(record_text)
    if row_codes and record_codes and not row_codes & record_codes:
        return -1
    overlap = len(tokenize(row_text) & tokenize(record_text))
    return overlap + (3 if row_codes & record_codes else 0)


def validate_insight_against_records(row: pd.Series | dict, records: list[dict[str, Any]]) -> tuple[str, str]:
    """Return (metrics_note, validation_verdict) for a single insight."""
    if not records:
        return (
            "No comparable metric table was extracted from the uploaded files.",
            "No cross-source quantitative validation available.",
        )

    get = row.get
    insight = str(get("insight_text", "") or "")
    row_text = f"{get('theme', '') or ''} {insight}"
    normalised_insight = normalize_compare_text(insight)
    direction = claim_direction(insight)

    best: tuple[int, str, str] | None = None
    fallback: tuple[str, str] | None = None

    for record in records:
        score = _relevance(row_text, record)
        if score < 2:
            continue
        values = [(v["label"], v.get("parent", ""), float(v["value"])) for v in record["values"] if v.get("value") is not None]
        if len(values) < 2:
            continue

        claimed = [
            item for item in values
            if len(normalize_compare_text(item[0])) >= 4 and normalize_compare_text(item[0]) in normalised_insight
        ]

        ranked = sorted(values, key=lambda item: item[2], reverse=True)
        top_label, _, top_value = ranked[0]
        bottom_label, _, bottom_value = ranked[-1]
        comparison = "; ".join(f"{label}={format_metric_value(value)}" for label, _, value in ranked[:6])
        base_note = f" Base: {record['base_note']}." if record.get("base_note") else ""
        metrics_note = (
            f"{record.get('title') or 'Matched table'} / {record.get('metric')}: {comparison}; "
            f"highest={top_label} ({format_metric_value(top_value)}); lowest={bottom_label} "
            f"({format_metric_value(bottom_value)}).{base_note} "
            f"Source={record.get('source_file')}, sheet={record.get('sheet')}."
        )

        if not claimed:
            if fallback is None:
                fallback = (metrics_note, "A related table was found, but the option named in the insight "
                                          "does not match any row or column label (no exact claimed option).")
            continue

        claimed_label, claimed_parent, claimed_value = claimed[0]
        peers = [item for item in values if item[1] == claimed_parent] if claimed_parent else values
        if len(peers) < 2:
            peers = values
        peer_ranked = sorted(peers, key=lambda item: item[2], reverse=True)
        true_top, true_bottom = peer_ranked[0], peer_ranked[-1]

        if direction == "high" and normalize_compare_text(claimed_label) != normalize_compare_text(true_top[0]):
            verdict = (
                f"Contradiction: the insight names {claimed_label} as highest, but the source shows "
                f"{claimed_label}={format_metric_value(claimed_value)} while {true_top[0]} is higher at "
                f"{format_metric_value(true_top[2])}."
            )
        elif direction == "low" and normalize_compare_text(claimed_label) != normalize_compare_text(true_bottom[0]):
            verdict = (
                f"Contradiction: the insight names {claimed_label} as lowest, but the source shows "
                f"{claimed_label}={format_metric_value(claimed_value)} while {true_bottom[0]} is lower at "
                f"{format_metric_value(true_bottom[2])}."
            )
        else:
            verdict = (
                f"Supported: {claimed_label}={format_metric_value(claimed_value)} in the matched source "
                f"(highest={true_top[0]} {format_metric_value(true_top[2])})."
            )

        candidate = (score, metrics_note, verdict)
        if best is None or candidate[0] > best[0]:
            best = candidate

    if best:
        return best[1], best[2]
    if fallback:
        return fallback
    return (
        "No matching metric was found for this insight in the uploaded tables.",
        "No directly matching quantitative metric found; the researcher should check the source tables manually.",
    )
