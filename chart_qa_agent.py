from __future__ import annotations

import io
import re
from typing import Any

import pandas as pd
from pptx import Presentation


CHART_QA_COLUMNS = [
    "chart_validation_status",
    "chart_metric_used",
    "chart_series_used",
    "chart_categories_compared",
    "claimed_item_value",
    "true_highest_or_lowest_item",
    "chart_validation_reason",
]

STOPWORDS = {
    "the", "and", "for", "with", "that", "this", "from", "into", "each",
    "study", "research", "insight", "insights", "report", "data", "analysis",
    "output", "client", "business", "market", "are", "than", "under", "over",
    "most", "highest", "lowest", "least", "strongest", "weakest", "better",
    "worse", "driving", "drives", "variant", "variants", "feature", "benefit",
}

POSITIVE_CLAIM_CUES = [
    "most", "highest", "strongest", "best", "winner", "winning", "leads",
    "higher than any", "more than any", "outperforms", "top", "most noticeable",
    "most recalled", "stronger recall", "higher recall", "most effective",
    "driving higher", "highest scoring", "strongest performer",
]

NEGATIVE_CLAIM_CUES = [
    "least", "lowest", "weakest", "poorest", "worst", "lower than any",
    "underperforms", "least noticeable", "least recalled", "least effective",
]

NEGATIVE_METRIC_CUES = [
    "no", "not sure", "cant recall", "can t recall", "low rating", "low ratings",
    "detractor", "detractors", "negative", "complaint", "complaints",
    "not satisfied", "disagree", "reject", "unaware", "not aware",
    "not noticed", "did not notice",
]

POSITIVE_METRIC_CUES = [
    "yes", "top box", "satisfied", "promoter", "purchase intent", "appeal",
    "recall", "noticed", "notice", "aware", "consider", "like", "prefer",
    "2-3 times", "4-5 times", "more than 5 times", "agree",
]


def compact_text(value: object, limit: int = 360) -> str:
    text = " ".join(str(value or "").replace("\n", " ").split()).strip()
    if len(text) <= limit:
        return text
    return text[: limit - 3].rstrip() + "..."


def normalize_text(value: object) -> str:
    text = str(value or "").lower().replace("&", " and ")
    text = re.sub(r"[^a-z0-9]+", " ", text)
    return " ".join(text.split())


def tokenize(text: object) -> set[str]:
    return {
        token
        for token in re.findall(r"[a-zA-Z][a-zA-Z0-9-]{2,}", str(text or "").lower())
        if token not in STOPWORDS
    }


def question_codes(text: object) -> set[str]:
    return set(re.findall(r"\b[a-z]{1,4}\d+[a-z0-9_]*\b", str(text or "").lower()))


def value_as_number(value: object) -> float | None:
    if value is None or str(value).strip() == "":
        return None
    if isinstance(value, (int, float)) and not pd.isna(value):
        return float(value)
    text = str(value).strip().replace(",", "")
    pct = "%" in text
    text = text.replace("%", "")
    try:
        number = float(text)
    except ValueError:
        return None
    return number / 100 if pct else number


def format_metric_value(value: float) -> str:
    if -1 <= value <= 1:
        return f"{value * 100:.0f}%"
    if float(value).is_integer():
        return f"{value:.0f}"
    return f"{value:.1f}"


def file_bytes(uploaded_file) -> bytes:
    if uploaded_file is None:
        return b""
    if hasattr(uploaded_file, "getvalue"):
        return uploaded_file.getvalue()
    pos = uploaded_file.tell() if hasattr(uploaded_file, "tell") else None
    if hasattr(uploaded_file, "seek"):
        uploaded_file.seek(0)
    data = uploaded_file.read()
    if pos is not None and hasattr(uploaded_file, "seek"):
        uploaded_file.seek(pos)
    return data


def slide_title_from_shapes(slide) -> str:
    for shape in slide.shapes:
        text = str(getattr(shape, "text", "") or "").replace("\n", " ").strip()
        if text and 2 <= len(text.split()) <= 18:
            return compact_text(text, 120)
    return "PowerPoint slide"


def has_positive_decisive_claim(text: object) -> bool:
    lowered = str(text or "").lower()
    return any(cue in lowered for cue in POSITIVE_CLAIM_CUES) and not any(
        cue in lowered for cue in NEGATIVE_CLAIM_CUES
    )


def has_negative_decisive_claim(text: object) -> bool:
    lowered = str(text or "").lower()
    return any(cue in lowered for cue in NEGATIVE_CLAIM_CUES)


def negative_metric_name(text: object) -> bool:
    normalized = normalize_text(text)
    return normalized in {"no", "not sure", "cant recall", "can t recall"} or any(
        cue in normalized for cue in NEGATIVE_METRIC_CUES
    )


def positive_metric_name(text: object) -> bool:
    normalized = normalize_text(text)
    return any(cue in normalized for cue in POSITIVE_METRIC_CUES)


def extract_ppt_chart_records(uploaded_files) -> list[dict[str, Any]]:
    records: list[dict[str, Any]] = []
    for uploaded_file in uploaded_files or []:
        source_name = getattr(uploaded_file, "name", "PowerPoint report")
        if not str(source_name).lower().endswith(".pptx"):
            continue
        data = file_bytes(uploaded_file)
        if not data:
            continue
        try:
            prs = Presentation(io.BytesIO(data))
        except Exception:
            continue

        for slide_index, slide in enumerate(prs.slides, start=1):
            slide_title = slide_title_from_shapes(slide)
            for shape in slide.shapes:
                if getattr(shape, "has_chart", False):
                    records.extend(
                        chart_records_from_shape(shape, source_name, slide_index, slide_title)
                    )
                if getattr(shape, "has_table", False):
                    records.extend(
                        table_records_from_shape(shape, source_name, slide_index, slide_title)
                    )
    return records


def chart_records_from_shape(shape, source_name: str, slide_index: int, slide_title: str) -> list[dict[str, Any]]:
    records: list[dict[str, Any]] = []
    chart = shape.chart
    try:
        categories = [str(category).replace("\n", " ").strip() for category in chart.plots[0].categories]
    except Exception:
        categories = []
    if not categories:
        return records
    chart_title = ""
    try:
        if chart.has_title:
            chart_title = str(chart.chart_title.text_frame.text or "").replace("\n", " ").strip()
    except Exception:
        chart_title = ""
    title = compact_text(f"PowerPoint slide {slide_index} {slide_title} {chart_title}".strip(), 220)
    for series in chart.series:
        try:
            series_name = str(getattr(series, "name", "") or "").replace("\n", " ").strip() or "Chart value"
            values = list(series.values)
        except Exception:
            continue
        metric_values = []
        for label, value in zip(categories, values):
            number = value_as_number(value)
            if number is None or not str(label).strip():
                continue
            metric_values.append({"label": str(label).strip(), "parent": "PowerPoint chart", "value": number})
        if len(metric_values) >= 2:
            records.append(
                {
                    "source_file": source_name,
                    "sheet": f"Slide {slide_index}",
                    "slide_number": slide_index,
                    "title": title,
                    "metric": series_name,
                    "values": metric_values,
                    "source_kind": "editable PowerPoint chart",
                }
            )
    return records


def table_records_from_shape(shape, source_name: str, slide_index: int, slide_title: str) -> list[dict[str, Any]]:
    records: list[dict[str, Any]] = []
    try:
        table = shape.table
    except Exception:
        return records
    rows = []
    for row in table.rows:
        rows.append([cell.text.replace("\n", " ").strip() for cell in row.cells])
    if len(rows) < 2:
        return records

    header = rows[0]
    for row_idx, row in enumerate(rows[1:], start=1):
        if not row:
            continue
        metric = row[0] or f"PowerPoint table row {row_idx}"
        values = []
        for col_idx, raw_value in enumerate(row[1:], start=1):
            label = header[col_idx] if col_idx < len(header) else f"Column {col_idx + 1}"
            number = value_as_number(raw_value)
            if number is None or not str(label).strip():
                continue
            values.append({"label": str(label).strip(), "parent": "PowerPoint table", "value": number})
        if len(values) >= 2:
            records.append(
                {
                    "source_file": source_name,
                    "sheet": f"Slide {slide_index}",
                    "slide_number": slide_index,
                    "title": compact_text(f"PowerPoint slide {slide_index} {slide_title}", 220),
                    "metric": metric,
                    "values": values,
                    "source_kind": "PowerPoint table",
                }
            )
    return records


def preferred_slide_number(row: dict[str, Any]) -> int | None:
    for field in ["insight_id", "source_slide", "source_file"]:
        match = re.search(r"PPT-S0*(\d+)|slide\s*0*(\d+)", str(row.get(field, "") or ""), re.I)
        if match:
            return int(match.group(1) or match.group(2))
    return None


def claimed_values_from_record(insight_text: str, record: dict[str, Any]) -> list[tuple[str, str, float]]:
    normalized_insight = normalize_text(insight_text)
    claimed = []
    for item in record.get("values", []):
        label = str(item.get("label", "") or "")
        normalized_label = normalize_text(label)
        if len(normalized_label) >= 4 and normalized_label in normalized_insight:
            claimed.append((label, str(item.get("parent", "") or ""), float(item.get("value"))))
    return claimed


def record_relevance(row: dict[str, Any], record: dict[str, Any]) -> int:
    text = f"{row.get('theme', '')} {row.get('insight_text', '')}"
    record_text = f"{record.get('title', '')} {record.get('metric', '')}"
    score = len(tokenize(text).intersection(tokenize(record_text)))
    row_codes = question_codes(text)
    record_codes = question_codes(record_text)
    if row_codes and record_codes and row_codes.intersection(record_codes):
        score += 4
    slide_no = preferred_slide_number(row)
    if slide_no and slide_no == record.get("slide_number"):
        score += 25
    metric_norm = normalize_text(record.get("metric", ""))
    insight_norm = normalize_text(row.get("insight_text", ""))
    if len(metric_norm) >= 4 and metric_norm in insight_norm:
        score += 4
    if has_positive_decisive_claim(text) and positive_metric_name(metric_norm):
        score += 3
    if has_positive_decisive_claim(text) and negative_metric_name(metric_norm):
        score -= 12
    return score


def compare_claim_to_record(row: dict[str, Any], record: dict[str, Any]) -> dict[str, str] | None:
    insight_text = str(row.get("insight_text", "") or "")
    positive_claim = has_positive_decisive_claim(insight_text)
    negative_claim = has_negative_decisive_claim(insight_text)
    if not positive_claim and not negative_claim:
        return None

    claimed = claimed_values_from_record(insight_text, record)
    if not claimed:
        return None

    claimed_label, claimed_parent, claimed_value = claimed[0]
    matched_values = [
        (str(item.get("label", "")), str(item.get("parent", "")), float(item.get("value")))
        for item in record.get("values", [])
        if item.get("value") is not None
    ]
    comparison_values = matched_values
    if claimed_parent:
        siblings = [item for item in matched_values if item[1] == claimed_parent]
        if len(siblings) >= 2:
            comparison_values = siblings
    if len(comparison_values) < 2:
        return None

    sorted_desc = sorted(comparison_values, key=lambda item: item[2], reverse=True)
    sorted_asc = sorted(comparison_values, key=lambda item: item[2])
    target_label, _, target_value = sorted_asc[0] if negative_claim else sorted_desc[0]
    target_word = "lowest" if negative_claim else "highest"
    categories = "; ".join(
        f"{label}={format_metric_value(value)}"
        for label, _, value in (sorted_asc[:5] if negative_claim else sorted_desc[:5])
    )
    metric = str(record.get("metric", "") or "")
    metric_note = (
        f"{record.get('title', 'Matched PowerPoint chart')} / {metric}: {categories}; "
        f"{target_word}={target_label} ({format_metric_value(target_value)}). "
        f"Source={record.get('source_file', '')}, {record.get('sheet', '')}."
    )

    status = "PASS"
    if normalize_text(claimed_label) != normalize_text(target_label):
        status = "FAIL"
        reason = (
            f"Chart contradiction: insight claims {claimed_label} is {target_word}/most decisive, "
            f"but the chart shows {claimed_label}={format_metric_value(claimed_value)} and "
            f"{target_label} is {target_word} at {format_metric_value(target_value)}."
        )
    elif positive_claim and negative_metric_name(metric):
        status = "FAIL"
        reason = (
            f"Opposite-condition contradiction: the claim is positive, but it is being checked against "
            f"the negative/opposite chart series '{metric}'. Validate against the positive KPI before reporting."
        )
    else:
        reason = (
            f"Chart support: claimed option {claimed_label} matches the {target_word} chart value "
            f"for '{metric}' ({format_metric_value(claimed_value)})."
        )

    return {
        "chart_validation_status": status,
        "chart_metric_used": compact_text(str(record.get("title", "")), 280),
        "chart_series_used": metric,
        "chart_categories_compared": compact_text(categories, 900),
        "claimed_item_value": f"{claimed_label}={format_metric_value(claimed_value)}",
        "true_highest_or_lowest_item": f"{target_label}={format_metric_value(target_value)}",
        "chart_validation_reason": reason,
        "quantitative_metrics_validated": metric_note,
        "cross_source_validation": reason if status == "FAIL" else reason.replace("Chart support:", "Supported:"),
    }


def empty_chart_result(reason: str = "No direct chart claim found for chart-level validation.") -> dict[str, str]:
    return {
        "chart_validation_status": "NOT_APPLICABLE",
        "chart_metric_used": "",
        "chart_series_used": "",
        "chart_categories_compared": "",
        "claimed_item_value": "",
        "true_highest_or_lowest_item": "",
        "chart_validation_reason": reason,
        "quantitative_metrics_validated": "",
        "cross_source_validation": "",
    }


def validate_chart_claim_against_records(row: dict[str, Any], records: list[dict[str, Any]]) -> dict[str, str]:
    if not records:
        return empty_chart_result("No machine-readable PowerPoint chart/table data found.")

    insight_text = str(row.get("insight_text", "") or "")
    if not has_positive_decisive_claim(insight_text) and not has_negative_decisive_claim(insight_text):
        return empty_chart_result()

    preferred_slide = preferred_slide_number(row)
    candidate_records = records
    if preferred_slide:
        same_slide_records = [record for record in records if record.get("slide_number") == preferred_slide]
        if same_slide_records:
            candidate_records = same_slide_records

    candidates = []
    for record in candidate_records:
        compared = compare_claim_to_record(row, record)
        if not compared:
            continue
        candidates.append((record_relevance(row, record), compared))

    if not candidates:
        return empty_chart_result(
            "Decisive chart claim detected, but the claimed option could not be matched to a readable chart category."
        )
    return sorted(candidates, key=lambda item: item[0], reverse=True)[0][1]
